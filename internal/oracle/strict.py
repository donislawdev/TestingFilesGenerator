"""Structural checks written against the format specifications.

These exist because the real readers are tolerant on purpose. Measured:
ffprobe ignores a wrong size in the RIFF header, Pillow does not verify the
checksum of an ancillary PNG chunk, and a PDF reader rebuilds a cross
reference table whose offset is off by one. All three are genuine defects that
a tolerant reader hides.

So this is a second layer: an independent implementation, in another language,
written to the specification rather than to what a viewer will put up with. It
complements the tolerant readers rather than replacing them - a file has to
pass both.

Usage: python strict.py <format> <path>
Prints "OK <detail>" or "FAIL <reason>".
"""

import re
import struct
import sys
import zlib


def fail(msg):
    print("FAIL " + msg)
    sys.exit(1)


def ok(msg):
    print("OK " + msg)
    sys.exit(0)


def check_png(data):
    if data[:8] != b"\x89PNG\r\n\x1a\n":
        fail("the signature is not a PNG signature")

    pos, seen, chunks = 8, [], 0
    while pos < len(data):
        if pos + 8 > len(data):
            fail(f"a chunk header runs past the end of the file at offset {pos}")
        length = struct.unpack(">I", data[pos:pos + 4])[0]
        kind = data[pos + 4:pos + 8]
        body = data[pos + 8:pos + 8 + length]
        if len(body) != length:
            fail(f"chunk {kind!r} declares {length} B and the file holds {len(body)}")
        want = struct.unpack(">I", data[pos + 8 + length:pos + 12 + length])[0]
        got = zlib.crc32(kind + body) & 0xFFFFFFFF
        if want != got:
            fail(f"chunk {kind.decode('latin1')} has checksum {want:08x} and its bytes give {got:08x}")
        seen.append(kind)
        chunks += 1
        pos += 12 + length

    if seen[0] != b"IHDR":
        fail("the first chunk is not IHDR")
    if seen[-1] != b"IEND":
        fail("the last chunk is not IEND")
    if b"IDAT" not in seen:
        fail("there is no image data")
    ok(f"{chunks} chunks, every checksum matches")


def check_wav(data):
    if data[:4] != b"RIFF" or data[8:12] != b"WAVE":
        fail("this is not a RIFF WAVE file")

    declared = struct.unpack("<I", data[4:8])[0]
    actual = len(data) - 8
    if declared != actual:
        fail(f"the RIFF header declares {declared} B after itself and the file holds {actual}")

    pos, seen, total = 12, [], 0
    while pos + 8 <= len(data):
        kind = data[pos:pos + 4]
        length = struct.unpack("<I", data[pos + 4:pos + 8])[0]
        body_end = pos + 8 + length
        if body_end > len(data):
            fail(f"chunk {kind!r} declares {length} B and runs past the end")
        seen.append(kind.decode("latin1").strip())
        total += 1
        pos = body_end
        # A chunk of odd length is padded to an even boundary, except when it
        # is the last one in the file.
        if length % 2 and pos < len(data):
            if data[pos:pos + 1] != b"\x00":
                fail(f"chunk {kind!r} has an odd length and is not followed by an alignment byte")
            pos += 1

    if pos != len(data):
        fail(f"the chunks end at {pos} and the file is {len(data)} B")
    if "fmt" not in seen:
        fail("there is no format chunk")
    if "data" not in seen:
        fail("there is no data chunk")
    ok(f"{total} chunks: {', '.join(seen)}")


def check_pdf(data):
    if not data.startswith(b"%PDF-"):
        fail("the file does not start with a PDF header")
    if b"%%EOF" not in data[-64:]:
        fail("the end of file marker is not near the end")

    idx = data.rfind(b"startxref")
    if idx < 0:
        fail("there is no startxref keyword")
    tail = data[idx + len(b"startxref"):].strip().split()[0]
    try:
        offset = int(tail)
    except ValueError:
        fail(f"startxref is followed by {tail!r} rather than a number")

    if offset >= len(data):
        fail(f"startxref points at {offset} and the file is {len(data)} B")
    if data[offset:offset + 4] != b"xref":
        near = data[max(0, offset - 8):offset + 12]
        fail(f"startxref points at {offset}, where the file holds {near!r} rather than the xref table")

    # Every offset in the table has to land on the object it claims.
    body = data[offset:]
    lines = body.split(b"\n")
    if len(lines) < 2:
        fail("the xref table has no entries")
    try:
        first, count = lines[1].split()
        count = int(count)
    except Exception:
        fail(f"the xref header is {lines[1]!r}")

    checked = 0
    for i in range(1, count):
        entry = lines[1 + i + 1] if 1 + i + 1 < len(lines) else b""
        parts = entry.split()
        if len(parts) < 3 or parts[2] != b"n":
            continue
        obj_at = int(parts[0])
        expected = str(i).encode() + b" 0 obj"
        if data[obj_at:obj_at + len(expected)] != expected:
            fail(f"the xref table sends object {i} to offset {obj_at}, "
                 f"where the file holds {data[obj_at:obj_at + 16]!r}")
        checked += 1
    ok(f"{checked} object offsets all land on their object")


def check_zip(data):
    end = data.rfind(b"PK\x05\x06")
    if end < 0:
        fail("there is no end of central directory record")
    comment_len = struct.unpack("<H", data[end + 20:end + 22])[0]
    if end + 22 + comment_len != len(data):
        fail(f"the comment declares {comment_len} B and {len(data) - end - 22} follow the record")
    entries = struct.unpack("<H", data[end + 10:end + 12])[0]
    cd_at = struct.unpack("<I", data[end + 16:end + 20])[0]
    if data[cd_at:cd_at + 4] != b"PK\x01\x02":
        fail(f"the central directory should start at {cd_at} and does not")
    ok(f"{entries} entries, comment {comment_len} B")


OCTET = r"(?:0|[1-9][0-9]{0,2})"
ADDR_V4 = OCTET + r"(?:\." + OCTET + r"){3}"
ADDR_V6 = r"[0-9a-f]{1,4}(?::[0-9a-f]{1,4}){7}"
ADDR = r"(?:" + ADDR_V4 + r"|" + ADDR_V6 + r")"
ISO_TIME = r"[0-9]{4}-[0-9]{2}-[0-9]{2}T[0-9]{2}:[0-9]{2}:[0-9]{2}\.[0-9]{6}[+-][0-9]{2}:[0-9]{2}"
REQUEST = (r" - - \[[^\]]+\] \"(?:GET|POST|PUT|DELETE|HEAD|PATCH) /\S* HTTP/1\.[01]\""
           r" [1-5][0-9]{2} [0-9]+")

# One pattern per shape this generator writes. Ordered most specific first,
# because nginx is combined plus one more quoted field and combined is common
# plus two, so a looser pattern tried first would claim a line it does not
# describe.
LOG_SHAPES = [
    ("nginx", re.compile(r"^" + ADDR + REQUEST + r" \"[^\"]*\" \"[^\"]*\" \"[^\"]*\"$")),
    ("apache-combined", re.compile(r"^" + ADDR + REQUEST + r" \"[^\"]*\" \"[^\"]*\"$")),
    ("apache-common", re.compile(r"^" + ADDR + REQUEST + r"$")),
    ("syslog", re.compile(r"^" + ISO_TIME + r" \S+ \S+\[[0-9]+\]: .*$")),
    ("plain", re.compile(r"^" + ISO_TIME + r" (?:DEBUG|INFO|WARN|ERROR) .*$")),
]


def check_log(data):
    """Every line is a whole entry, and every entry is the SAME shape.

    A log is read line by line, so a line that is not a whole entry is a broken
    file however right its length is. And "the last line is truncated" is what
    a real log looks like caught mid rotation, so the defect reads as realism
    unless something checks.

    Written to the format rather than to our output. The octet pattern refuses
    a leading zero: real logs write 93, not 093, and a leading zero is read as
    octal by some address parsers - where 069 is not even valid octal. Our
    generator produced padded octets until 2026-08-01, and a checker written to
    match it would have blessed that instead of catching it.

    Since 2026-08-31 the generator writes six shapes, and this checker is told
    which one only by the file - it is handed a format id and a path, never the
    recipe. So it takes the shape from the first entry and holds every other
    line to THAT one. Asking each line merely to be valid on its own would pass
    a file that changed shape half way down, which no reader could load.
    """
    try:
        text = data.decode("utf-8")
    except UnicodeDecodeError as exc:
        fail(f"not valid UTF-8: {exc}")

    if not text.endswith(("\n", "\r\n")):
        fail("the file does not end with a newline, so the last entry is unterminated")

    # Line endings have to be consistent. A file mixing them is the shape a
    # careless writer produces and a strict reader rejects.
    crlf = "\r\n" in text
    if crlf and re.search(r"(?<!\r)\n", text):
        fail("the file mixes CRLF and bare LF endings")
    lines = text.replace("\r\n", "\n").rstrip("\n").split("\n")

    entries, shape = 0, None
    for number, line in enumerate(lines, start=1):
        if line.startswith("# "):
            continue  # the label line, which says in words that it is not an entry
        if line.startswith('{"label":'):
            continue  # the same thing where a comment would not parse
        if shape is None:
            shape = detect_log_shape(line, number)
        check_log_line(line, number, shape)
        entries += 1

    if entries == 0:
        fail("the file holds no entries at all")
    ok(f"{entries} {shape} entries, all whole, {'crlf' if crlf else 'lf'} endings")


def detect_log_shape(line, number):
    if line.startswith("{"):
        return "json-lines"
    for name, pattern in LOG_SHAPES:
        if pattern.match(line):
            return name
    fail(f"line {number} is not a whole entry in any shape this tool writes: {line[:90]!r}")


def check_log_line(line, number, shape):
    if shape == "json-lines":
        # A real parser rather than a pattern, which is what makes this the one
        # shape here checked by somebody else's implementation.
        import json
        try:
            obj = json.loads(line)
        except ValueError as exc:
            fail(f"line {number} is not valid JSON: {exc}")
        if not isinstance(obj, dict):
            fail(f"line {number} is JSON but not an object: {line[:60]!r}")
        for key in ("time", "level", "msg"):
            if key not in obj:
                fail(f"line {number} has no {key!r} field")
        return

    pattern = dict(LOG_SHAPES)[shape]
    if not pattern.match(line):
        fail(f"line {number} is not a whole {shape} entry: {line[:90]!r}")

    if shape in ("nginx", "apache-combined", "apache-common"):
        # An octet above 255 parses as a number and is not an address.
        address = line.split(" ", 1)[0]
        if "." in address and any(int(part) > 255 for part in address.split(".")):
            fail(f"line {number} has an octet above 255: {address}")


CSV_DELIMITERS = {"comma": ",", "semicolon": ";", "tab": "\t", "pipe": "|"}
CSV_LINE_ENDINGS = {"lf": "\n", "crlf": "\r\n"}


def check_csv(data, settings=None):
    """Every row carries the same columns, written to RFC 4180 by hand.

    Not csv.reader. That module is the tolerant reader used as the reference
    tool beside this one, and it accepts a ragged row without a word - a row
    with a column missing is exactly the defect this has to catch, and the file
    would still be the right size and still parse.

    The field size ceiling below is the measurement of 2026-08-01 turned into a
    rule: the Python csv module at its default settings refuses a field above
    131 072 B. Padding pushed into one field instead of through rows would sail
    past every size and determinism guard and break the reader a tester has
    nearest.

    The dialect is TOLD rather than worked out. A checker that sniffed the
    separator would agree with a file that used the wrong one - it would split
    on whatever it found and report a tidy table either way - and whether the
    file uses the separator that was ordered is a question for a guard reading
    the manifest, not for this. Told, it can still catch the defect that
    matters here: a header and its rows disagreeing about the separator.
    """
    FIELD_CEILING = 131072

    settings = settings or {}
    sep = CSV_DELIMITERS[settings.get("delimiter", "comma")]
    eol = CSV_LINE_ENDINGS[settings.get("line_ending", "lf")]
    has_header = settings.get("header", "true") == "true"

    try:
        text = data.decode("utf-8")
    except UnicodeDecodeError as exc:
        fail(f"not valid UTF-8: {exc}")

    if not text.endswith(eol):
        fail(f"the file does not end with {settings.get('line_ending', 'lf')}, "
             "so the last row is unterminated")

    rows, field, row, quoted, i = [], [], [], False, 0
    while i < len(text):
        ch = text[i]
        if quoted:
            if ch == '"':
                if i + 1 < len(text) and text[i + 1] == '"':
                    field.append('"')
                    i += 2
                    continue
                quoted = False
            else:
                field.append(ch)
            i += 1
            continue
        if ch == '"':
            if field:
                fail(f"row {len(rows) + 1} opens a quote in the middle of a field")
            quoted = True
        elif ch == sep:
            row.append("".join(field))
            field = []
        elif text.startswith(eol, i):
            row.append("".join(field))
            field = []
            rows.append(row)
            row = []
            i += len(eol)
            continue
        elif ch in "\r\n":
            # A bare terminator where the dialect says there should be another
            # one. Caught here rather than swept into a field, because a lone CR
            # inside a CRLF file is the shape a half converted writer produces
            # and it splits rows for some readers and not others.
            fail(f"row {len(rows) + 1} carries a bare "
                 f"{'CR' if ch == chr(13) else 'LF'} where the rows end with "
                 f"{settings.get('line_ending', 'lf')}")
        else:
            field.append(ch)
        i += 1

    if quoted:
        fail("a quoted field is never closed")
    if field or row:
        fail("the file ends in the middle of a row")
    # With a header there has to be something under it. Without one, a single
    # row is the whole file and is legal - there is simply nothing to compare
    # it against, which is a limit of this check rather than a fault in it.
    least = 2 if has_header else 1
    if len(rows) < least:
        fail(f"the table holds {len(rows)} row(s), so there is no data to check")

    columns = len(rows[0])
    if columns < 2:
        fail(f"the first row has {columns} column(s), so nothing is separated by "
             f"the {settings.get('delimiter', 'comma')} this file is meant to use")
    declares = "the header" if has_header else "the first row"
    for number, r in enumerate(rows, start=1):
        if len(r) != columns:
            fail(f"row {number} has {len(r)} fields and {declares} declares {columns}")
        for value in r:
            if len(value.encode("utf-8")) > FIELD_CEILING:
                fail(f"row {number} has a field of {len(value.encode('utf-8'))} B - "
                     f"the default Python reader refuses anything above {FIELD_CEILING} B")

    data = len(rows) - 1 if has_header else len(rows)
    ok(f"{data} data rows, {columns} columns each")


def check_json(data):
    """The document is an array of records, each on its own line.

    Parsing is CPython's json here and V8's parser in the reference tool beside
    it, which are two implementations in two languages. What this adds is the
    shape: that the file is records rather than one enormous value, and that it
    is laid out the way the manifest says it is.

    The value types come from the format document, not from what the generator
    happens to emit. A generator that quietly stopped writing booleans would be
    the right size and would still parse.
    """
    import json

    try:
        text = data.decode("utf-8")
    except UnicodeDecodeError as exc:
        fail(f"not valid UTF-8: {exc}")

    try:
        doc = json.loads(text)
    except ValueError as exc:
        fail(f"not valid JSON: {exc}")

    if not isinstance(doc, list):
        fail(f"the root is {type(doc).__name__} rather than an array")
    if not doc:
        fail("the array is empty")

    # The manifest states one record per line, so the line count has to match.
    lines = text.rstrip("\n").split("\n")
    if len(lines) != len(doc) + 2:
        fail(f"{len(doc)} records over {len(lines)} lines - the manifest says one record per line")

    kinds = set()
    keys = None
    for number, record in enumerate(doc, start=1):
        if not isinstance(record, dict):
            fail(f"record {number} is {type(record).__name__} rather than an object")
        if keys is None:
            keys = set(record)
        elif set(record) != keys:
            fail(f"record {number} has the keys {sorted(record)} and record 1 has {sorted(keys)}")
        for value in record.values():
            if value is None:
                kinds.add("null")
            elif isinstance(value, bool):
                kinds.add("bool")
            elif isinstance(value, (int, float)):
                kinds.add("number")
            elif isinstance(value, str):
                kinds.add("string")
            elif isinstance(value, list):
                kinds.add("array")
            elif isinstance(value, dict):
                kinds.add("object")

    wanted = {"null", "bool", "number", "string", "array", "object"}
    missing = wanted - kinds
    if missing:
        fail(f"no record carries a value of type {', '.join(sorted(missing))} - "
             "the point of a JSON fixture is that a parser meets every type")

    ok(f"{len(doc)} records, keys {sorted(keys)}")


def scan_xml(data):
    """Well formed, checked by hand against the XML specification.

    Deliberately not expat. That is the reference tool beside this one, so
    calling it here would run the same C parser twice and prove nothing the
    first run did not.

    What a hand scanner adds: tags balanced and correctly nested, a comment that
    never contains a double hyphen, and text where every ampersand starts a real
    entity reference. A raw ampersand is the classic way a generated document
    stops being well formed while staying exactly the right size.
    """
    try:
        text = data.decode("utf-8")
    except UnicodeDecodeError as exc:
        fail(f"not valid UTF-8: {exc}")

    if not text.startswith("<?xml "):
        fail("the document does not open with an XML declaration")

    entity = re.compile(r"&(?:amp|lt|gt|quot|apos|#[0-9]+|#x[0-9a-fA-F]+);")
    stack, elements, i, seen = [], 0, 0, {}

    while i < len(text):
        nxt = text.find("<", i)
        if nxt < 0:
            nxt = len(text)

        # Character data between two tags.
        chunk = text[i:nxt]
        for pos, ch in enumerate(chunk):
            if ch == "&" and not entity.match(chunk, pos):
                fail(f"a bare ampersand at offset {i + pos} - text has to escape it as &amp;")
        if chunk.strip() and not stack:
            fail(f"text outside the root element at offset {i}: {chunk.strip()[:40]!r}")
        if nxt == len(text):
            break

        if text.startswith("<!--", nxt):
            end = text.find("-->", nxt + 4)
            if end < 0:
                fail("a comment is never closed")
            if "--" in text[nxt + 4:end]:
                fail("a comment contains a double hyphen, which XML does not allow")
            i = end + 3
            continue

        if text.startswith("<?", nxt):
            end = text.find("?>", nxt + 2)
            if end < 0:
                fail("a processing instruction is never closed")
            i = end + 2
            continue

        end = text.find(">", nxt)
        if end < 0:
            fail(f"a tag opened at offset {nxt} is never closed")
        tag = text[nxt + 1:end]

        if tag.startswith("/"):
            name = tag[1:].strip()
            if not stack:
                fail(f"</{name}> closes an element that was never opened")
            opened = stack.pop()
            if opened != name:
                fail(f"</{name}> closes while <{opened}> is open")
        else:
            selfClosing = tag.endswith("/")
            body = tag[:-1] if selfClosing else tag
            name = body.split()[0] if body.split() else ""
            if not name:
                fail(f"an empty tag at offset {nxt}")
            # Every attribute value has to be quoted.
            for attr in re.finditer(r"(\w+)\s*=\s*(.?)", body[len(name):]):
                if attr.group(2) not in ('"', "'"):
                    fail(f"attribute {attr.group(1)} of <{name}> is not quoted")
            elements += 1
            seen[name] = seen.get(name, 0) + 1
            if not selfClosing:
                stack.append(name)
        i = end + 1

    if stack:
        fail(f"the document ends with {', '.join('<' + s + '>' for s in stack)} still open")
    return elements, seen


def check_xml(data):
    elements, _ = scan_xml(data)
    if elements < 2:
        fail(f"the document holds {elements} element(s), so there is nothing below the root")
    ok(f"{elements} elements, all balanced")


def check_svg(data):
    """An SVG that parses can still draw nothing.

    Inkscape stands beside this one and answers whether the file renders. What
    this adds is the shape of the document: the root really is an svg with a
    viewBox, and it holds drawable elements rather than only metadata. Written
    to the format, not to our output - the list below is what SVG calls a basic
    shape, and a generator that quietly stopped emitting them would be the right
    size and would still open.
    """
    elements, seen = scan_xml(data)

    text = data.decode("utf-8")
    if "<svg" not in text:
        fail("there is no svg root element")
    if 'xmlns="http://www.w3.org/2000/svg"' not in text:
        fail("the root carries no SVG namespace, so a renderer has no reason to draw it")
    if "viewBox=" not in text:
        fail("the root declares no viewBox")

    drawable = {"rect", "circle", "ellipse", "line", "polyline", "polygon", "path", "text"}
    shapes = sum(count for name, count in seen.items() if name in drawable)
    if shapes == 0:
        fail(f"the drawing holds {elements} element(s) and not one of them draws anything")
    ok(f"{shapes} drawable shapes out of {elements} elements")


def check_html(data):
    """Balanced, complete and with real blocks in the body.

    HTML is the weakest format in this project for checking, and that is a
    property of the format rather than of this machine. A parser is required to
    recover from almost anything, so the tolerant reader beside this one accepts
    documents nobody would want - "it parsed" is close to no information. This
    check is written to the specification instead: the void elements below are
    the HTML5 list, everything else has to close, and an ampersand has to start
    a real entity reference.
    """
    VOID = {"area", "base", "br", "col", "embed", "hr", "img", "input",
            "link", "meta", "param", "source", "track", "wbr"}
    BLOCKS = {"p", "h1", "h2", "h3", "ul", "ol", "table", "blockquote"}

    try:
        text = data.decode("utf-8")
    except UnicodeDecodeError as exc:
        fail(f"not valid UTF-8: {exc}")

    if not text.lower().startswith("<!doctype html>"):
        fail("the document does not open with an HTML5 doctype")
    if not text.rstrip().endswith("</html>"):
        fail("the document does not end with a closing html tag")

    entity = re.compile(r"&(?:[a-zA-Z][a-zA-Z0-9]{1,31}|#[0-9]+|#x[0-9a-fA-F]+);")
    tag = re.compile(r"<(/?)([a-zA-Z][a-zA-Z0-9]*)([^>]*)>")

    stack, blocks, i = [], 0, 0
    for m in tag.finditer(text):
        # Character data between the previous tag and this one.
        chunk = text[i:m.start()]
        for pos, ch in enumerate(chunk):
            if ch == "&" and not entity.match(chunk, pos):
                fail(f"a bare ampersand at offset {i + pos} - text has to escape it as &amp;")
        i = m.end()

        closing, name, rest = m.group(1) == "/", m.group(2).lower(), m.group(3)
        if name in VOID or rest.rstrip().endswith("/"):
            continue
        if closing:
            if not stack:
                fail(f"</{name}> closes an element that was never opened")
            opened = stack.pop()
            if opened != name:
                fail(f"</{name}> closes while <{opened}> is open")
        else:
            stack.append(name)
            if name in BLOCKS:
                blocks += 1

    if stack:
        fail(f"the document ends with {', '.join('<' + s + '>' for s in stack)} still open")
    if blocks == 0:
        fail("the body holds no block elements, so the page renders as nothing")
    ok(f"{blocks} blocks, all tags balanced")


def gzip_header_end(data):
    """Where the deflate stream starts, per RFC 1952 section 2.3.

    Walked by hand rather than handed to the gzip module, because the module
    skips these fields without checking that they are terminated inside the
    file - which is exactly the class of defect this layer exists to catch.
    Returns the offset and the header comment.
    """
    if data[:2] != b"\x1f\x8b":
        fail("there is no gzip magic at the start of the file")
    if data[2] != 8:
        fail(f"the compression method is {data[2]} and deflate is 8")
    flg = data[3]
    if flg & 0xE0:
        fail(f"reserved bits are set in the flag byte ({flg:#04x})")

    at = 10
    if flg & 0x04:
        if at + 2 > len(data):
            fail("the extra field length runs past the end of the file")
        xlen = struct.unpack("<H", data[at:at + 2])[0]
        at += 2 + xlen
    if flg & 0x08:
        end = data.find(b"\x00", at)
        if end < 0:
            fail("the original name field is not terminated")
        at = end + 1
    comment = b""
    if flg & 0x10:
        end = data.find(b"\x00", at)
        if end < 0:
            fail("the header comment is not terminated")
        comment = data[at:end]
        at = end + 1
    if flg & 0x02:
        at += 2
    if at > len(data):
        fail("the gzip header runs past the end of the file")
    return at, comment


def tar_entries(tar):
    """Every ustar header, with its stored checksum verified.

    The checksum is the reason this is here. Readers accept a header whose
    checksum is wrong - GNU tar warns and carries on - so a header corrupted in
    a way that leaves the sizes intact passes a listing and fails here.
    """
    if len(tar) % 512:
        fail(f"the tar stream is {len(tar)} B, which is not whole 512 B blocks")
    if len(tar) < 1024 or tar[-1024:] != b"\x00" * 1024:
        fail("the tar stream does not end with two zero blocks")

    entries = []
    at = 0
    limit = len(tar) - 1024
    while at < limit:
        head = tar[at:at + 512]
        if head[257:263] != b"ustar\x00":
            fail(f"the header at {at} does not carry the ustar magic")

        stored = head[148:156].split(b"\x00")[0].split(b" ")[0]
        try:
            claimed = int(stored, 8)
        except ValueError:
            fail(f"the header at {at} has an unreadable checksum field")
        # The checksum is computed with its own eight bytes read as spaces.
        actual = sum(head[:148]) + sum(head[156:]) + 8 * 32
        if claimed != actual:
            fail(f"the header at {at} claims checksum {claimed} and its bytes add up to {actual}")

        name = head[:100].split(b"\x00")[0].decode("ascii", "replace")
        raw = head[124:136].split(b"\x00")[0].strip()
        try:
            size = int(raw, 8) if raw else 0
        except ValueError:
            fail(f"the entry {name} has an unreadable size field")
        entries.append((name, size))
        at += 512 + (size + 511) // 512 * 512

    if at != limit:
        fail(f"the entries end at {at} and the end of archive marker starts at {limit}")
    return entries


def check_targz(data):
    at, comment = gzip_header_end(data)

    stream = zlib.decompressobj(-zlib.MAX_WBITS)
    try:
        tar = stream.decompress(data[at:]) + stream.flush()
    except zlib.error as exc:
        fail(f"the deflate stream does not decompress: {exc}")

    # Measured 2026-08-04: four readers out of five refuse a file with anything
    # after the stream, so the trailer has to be the last eight bytes and there
    # has to be nothing behind it.
    trailer = stream.unused_data
    if len(trailer) != 8:
        fail(f"exactly 8 B of trailer should follow the stream and {len(trailer)} B do")
    crc, isize = struct.unpack("<II", trailer)
    if crc != zlib.crc32(tar) & 0xFFFFFFFF:
        fail("the CRC32 in the trailer does not match the data it covers")
    if isize != len(tar) & 0xFFFFFFFF:
        fail(f"the trailer says {isize} B uncompressed and the stream holds {len(tar)}")

    entries = tar_entries(tar)
    if b"\x00" in comment:
        fail("the header comment holds a zero byte, which ends the field early")
    ok(f"{len(entries)} entries, comment {len(comment)} B, tar {len(tar)} B")


def check_bmp(data):
    if data[:2] != b"BM":
        fail("the signature is not a BMP signature")
    if len(data) < 54:
        fail(f"the file is {len(data)} B and the two headers alone are 54 B")

    declared, offbits = struct.unpack("<I", data[2:6])[0], struct.unpack("<I", data[10:14])[0]
    # A tolerant reader never looks at this. A wrong size here is exactly the
    # kind of defect that passes every viewer and breaks a parser.
    if declared != len(data):
        fail(f"the header says the file is {declared} B and it is {len(data)} B")

    header = struct.unpack("<I", data[14:18])[0]
    if header != 40:
        fail(f"the info header is {header} B, and this generator writes the 40 B form")

    width, height = struct.unpack("<ii", data[18:26])
    planes, depth = struct.unpack("<HH", data[26:30])
    compression, image_bytes = struct.unpack("<II", data[30:38])

    if width < 1:
        fail(f"the width is {width}")
    if height == 0:
        fail("the height is zero")
    if planes != 1:
        fail(f"the header declares {planes} colour planes and a BMP has one")
    if depth not in (1, 4, 8, 16, 24, 32):
        fail(f"the header declares {depth} bits per pixel")
    if compression != 0:
        fail(f"the header declares compression {compression} and this generator writes none")

    rows = abs(height)
    stride = (width * depth + 31) // 32 * 4
    if image_bytes not in (0, stride * rows):
        fail(f"the header says the pixels are {image_bytes} B and the geometry gives {stride * rows} B")
    if offbits < 54:
        fail(f"the pixels are said to start at {offbits}, inside the headers")
    if offbits + stride * rows > len(data):
        fail(f"the pixels run to {offbits + stride * rows} and the file ends at {len(data)}")

    gap = offbits - 54
    ok(f"{width}x{rows}, {depth} bit, stride {stride} B, gap {gap} B, size field agrees")


def gif_blocks(data, pos):
    """Walk a chain of sub blocks and return where it ends."""
    while True:
        if pos >= len(data):
            fail("a chain of sub blocks runs off the end of the file")
        size = data[pos]
        pos += 1
        if size == 0:
            return pos
        if pos + size > len(data):
            fail(f"a sub block of {size} B runs past the end of the file")
        pos += size


def check_gif(data):
    if data[:6] not in (b"GIF87a", b"GIF89a"):
        fail("the signature is not a GIF signature")
    if len(data) < 14:
        fail(f"the file is {len(data)} B and the screen descriptor alone needs 13 B")

    width, height, packed = struct.unpack("<HHB", data[6:11])
    pos = 13
    table = 0
    if packed & 0x80:
        table = 3 * (1 << ((packed & 0x07) + 1))
        pos += table
        if pos > len(data):
            fail(f"the colour table of {table} B runs past the end of the file")

    images, extensions, comment = 0, 0, 0
    while True:
        if pos >= len(data):
            fail("the file ends with no trailer")
        marker = data[pos]
        pos += 1
        if marker == 0x3B:
            break
        if marker == 0x21:
            if pos >= len(data):
                fail("an extension has no label")
            label = data[pos]
            pos += 1
            if label == 0xFE:
                start = pos
                pos = gif_blocks(data, pos)
                comment += pos - start - 1
            else:
                pos = gif_blocks(data, pos)
            extensions += 1
            continue
        if marker == 0x2C:
            if pos + 9 > len(data):
                fail("an image descriptor runs past the end of the file")
            local = data[pos + 8]
            pos += 9
            if local & 0x80:
                pos += 3 * (1 << ((local & 0x07) + 1))
            if pos >= len(data):
                fail("an image has no code size")
            pos += 1
            pos = gif_blocks(data, pos)
            images += 1
            continue
        fail(f"the byte {marker:#04x} at offset {pos - 1} is not a block marker")

    # A tolerant reader stops at the trailer and never looks further. Bytes
    # after it are outside the data stream, so a file carrying them is not the
    # file it claims to be.
    if pos != len(data):
        fail(f"the trailer is at {pos - 1} and the file holds {len(data) - pos} B after it")
    if images < 1:
        fail("there is no image block")
    ok(f"{width}x{height}, table {table} B, {images} image(s), "
       f"{extensions} extension(s), {comment} B of comment")


def check_ico(data):
    if len(data) < 6:
        fail(f"the file is {len(data)} B and the directory header alone is 6 B")
    reserved, kind, count = struct.unpack("<HHH", data[:6])
    if reserved != 0:
        fail(f"the reserved field is {reserved} and it has to be zero")
    if kind != 1:
        fail(f"the type is {kind} and an icon is 1")
    if count < 1:
        fail("the directory says it holds no images")
    if 6 + 16 * count > len(data):
        fail(f"the directory of {count} entries runs past the end of the file")

    described = []
    for i in range(count):
        at = 6 + 16 * i
        side_w, side_h = data[at], data[at + 1]
        length, offset = struct.unpack("<II", data[at + 8:at + 16])
        want_w = side_w if side_w else 256
        want_h = side_h if side_h else 256
        if offset < 6 + 16 * count:
            fail(f"image {i} starts at {offset}, inside the directory")
        if offset + length > len(data):
            fail(f"image {i} runs to {offset + length} and the file ends at {len(data)}")
        blob = data[offset:offset + length]
        described.append((want_w, want_h, len(blob), inside_icon(i, blob, want_w, want_h)))

    shapes = ", ".join(f"{w}x{h} {kindname} {n} B" for w, h, n, kindname in described)
    ok(f"{count} image(s): {shapes}")


def inside_icon(index, blob, want_w, want_h):
    """Say what sits inside one directory entry, and check it against the entry."""
    if blob[:8] == b"\x89PNG\r\n\x1a\n":
        if len(blob) < 24 or blob[12:16] != b"IHDR":
            fail(f"image {index} starts like a PNG and has no IHDR")
        w, h = struct.unpack(">II", blob[16:24])
        if (w, h) != (want_w, want_h):
            fail(f"the directory says image {index} is {want_w}x{want_h} and the PNG inside is {w}x{h}")
        return "png"

    if len(blob) < 40:
        fail(f"image {index} is {len(blob)} B, too small for a bitmap header")
    header = struct.unpack("<I", blob[:4])[0]
    if header != 40:
        fail(f"image {index} has a {header} B bitmap header, and 40 B is the form icons use")
    w, h = struct.unpack("<ii", blob[4:12])
    depth = struct.unpack("<H", blob[14:16])[0]
    # The height of an embedded bitmap counts the colour rows and the mask
    # rows together, so it is twice the height the directory states. Getting
    # this wrong gives an icon that draws upside down or half transparent, and
    # no reader complains about the header itself.
    if w != want_w:
        fail(f"the directory says image {index} is {want_w} wide and the bitmap says {w}")
    if h != want_h * 2:
        fail(f"the bitmap for image {index} declares height {h} and it has to be twice {want_h}")
    colour = (w * depth + 31) // 32 * 4 * want_h
    mask = (w + 31) // 32 * 4 * want_h
    if 40 + colour + mask != len(blob):
        fail(f"image {index} needs {40 + colour + mask} B for header, pixels and mask and holds {len(blob)}")
    return f"{depth} bit bitmap"


OPC_NS_CT = "{http://schemas.openxmlformats.org/package/2006/content-types}"
OPC_NS_REL = "{http://schemas.openxmlformats.org/package/2006/relationships}"

# Measured 2026-08-19: a reader stops accepting the package above this.
OPC_COMMENT_LIMIT = 512


def opc_open(data):
    import io as _io
    import zipfile
    try:
        return zipfile.ZipFile(_io.BytesIO(data))
    except Exception as exc:
        fail("the package is not a readable ZIP: %s" % exc)


def opc_check(data, roots, kind, key):
    """Checks every OOXML package has to pass, whatever is inside it.

    Written to the packaging specification rather than to what a viewer will
    put up with. Two of these were found the hard way and neither is visible to
    a ZIP tool: an entry carrying its sizes after the data instead of in the
    header, which LibreOffice's Writer and Impress filters refuse while its Calc
    filter accepts, and an archive comment above 512 B, which every ZIP reader
    accepts and no OOXML reader does.
    """
    import xml.etree.ElementTree as ET

    z = opc_open(data)
    names = [i.filename for i in z.infolist()]

    if len(z.comment) > OPC_COMMENT_LIMIT:
        fail("the archive comment is %d B and a reader stops at %d B"
             % (len(z.comment), OPC_COMMENT_LIMIT))

    for info in z.infolist():
        if info.flag_bits & 0x08:
            fail("the entry %s carries its sizes after the data instead of in "
                 "the header, which an Office reader refuses" % info.filename)

    if "[Content_Types].xml" not in names:
        fail("there is no [Content_Types].xml, so nothing says what the parts are")
    if "_rels/.rels" not in names:
        fail("there is no _rels/.rels, so nothing points at the document")

    try:
        types = ET.fromstring(z.read("[Content_Types].xml"))
    except ET.ParseError as exc:
        fail("[Content_Types].xml is not well formed: %s" % exc)

    defaults = {e.get("Extension", "").lower(): e.get("ContentType")
                for e in types.findall(OPC_NS_CT + "Default")}
    overrides = {e.get("PartName"): e.get("ContentType")
                 for e in types.findall(OPC_NS_CT + "Override")}

    for part, ct in overrides.items():
        if part.lstrip("/") not in names:
            fail("the content types name %s and the package does not hold it" % part)
        if not ct:
            fail("the content types give %s no type" % part)

    # Every part needs a type, from an override or from its extension. The
    # reader measured here does not check this - the specification does.
    for name in names:
        if name.endswith("/"):
            continue
        ext = name.rsplit(".", 1)[-1].lower() if "." in name else ""
        if ("/" + name) in overrides:
            continue
        if ext in defaults:
            continue
        fail("the part %s has no content type, by override or by extension" % name)

    roots_found = []
    for name in names:
        if not name.endswith(".xml") and not name.endswith(".rels"):
            continue
        try:
            root = ET.fromstring(z.read(name))
        except ET.ParseError as exc:
            fail("the part %s is not well formed XML: %s" % (name, exc))
        roots_found.append((name, root.tag))

    # Relationship targets have to exist, or a reader follows a link to nothing.
    for name in names:
        if not name.endswith(".rels"):
            continue
        base = name.rsplit("_rels/", 1)[0]
        rels = ET.fromstring(z.read(name))
        for rel in rels.findall(OPC_NS_REL + "Relationship"):
            if rel.get("TargetMode") == "External":
                continue
            target = rel.get("Target", "")
            resolved = opc_resolve(base, target)
            if resolved not in names:
                fail("%s points at %s and the package does not hold it"
                     % (name, resolved))

    for want in roots:
        if want not in names:
            fail("a %s needs %s and the package does not hold it" % (kind, want))

    padding = [i for i in z.infolist() if i.filename.startswith("tfg/")]
    pad = padding[0].file_size if padding else 0
    try:
        read = opc_readback(data, key)
    except SystemExit:
        raise
    except Exception as exc:
        # A reader refusing the package is a verdict about the package, and it
        # has to read as one. Letting the exception out gave a stack trace and
        # "the checker said nothing useful", which names the wrong culprit.
        fail("an independent reader would not read the package: %s: %s"
             % (type(exc).__name__, str(exc)[:120]))
    ok("%s, %d parts, %d B of padding, comment %d B, %s"
       % (kind, len(names), pad, len(z.comment), read))


def opc_resolve(base, target):
    """Turns a relationship target into a path inside the package."""
    import posixpath
    if target.startswith("/"):
        return target.lstrip("/")
    return posixpath.normpath(posixpath.join(base, target)).replace("\\", "/")



def opc_readback(data, key):
    """Czyta pakiet NIEZALEZNA implementacja OPC i sprawdza, ze cos w nim jest.

    Do 2026-08-19 czytnik OOXML na tej maszynie byl jeden - LibreOffice - i to
    bylo zapisane jako najslabszy pomiar w projekcie. Wlasciciel doinstalowal
    python-docx, openpyxl i python-pptx, czyli trzy osobne implementacje, zadna
    nie dzielaca kodu z LibreOffice ani ze soba.

    Stoi to tutaj, a nie wsrod wyrocznie, z powodu, ktory warto znac: format
    deklaruje JEDNA wyrocznie (`Descriptor.Oracle`), a to pole wychodzi do
    `tfg formats --json`, wiec poszerzenie go jest zmiana nosna. Warstwa
    strukturalna juz biegnie per format i pyta „czy plik jest naprawde dobrze
    zbudowany" - a biblioteka odmawiajaca odczytu odpowiada dokladnie na to.

    Brak biblioteki NIE jest porazka. Zostaje wtedy nazwany w wyniku, zeby
    przebieg, ktory sprawdzil mniej, nie wygladal jak ten, ktory sprawdzil
    wszystko.
    """
    import io as _io

    if key == "docx":
        try:
            from docx import Document
        except ImportError:
            return "no reader: python-docx"
        texts = [p.text for p in Document(_io.BytesIO(data)).paragraphs]
        if not texts:
            fail("python-docx otworzyl dokument i nie znalazl w nim ani jednego akapitu")
        return "python-docx odczytal %d akapit(ow)" % len(texts)

    if key == "xlsx":
        try:
            from openpyxl import load_workbook
        except ImportError:
            return "no reader: openpyxl"
        book = load_workbook(_io.BytesIO(data), read_only=True)
        sheet = book[book.sheetnames[0]]
        rows = sum(1 for _ in sheet.iter_rows(values_only=True))
        book.close()
        if rows < 1:
            fail("openpyxl otworzyl skoroszyt i nie znalazl w nim ani jednego wiersza")
        return "openpyxl odczytal %d wiersz(y)" % rows

    if key == "pptx":
        try:
            from pptx import Presentation
        except ImportError:
            return "no reader: python-pptx"
        deck = Presentation(_io.BytesIO(data))
        # Chodzimy po KSZTALTACH kazdego slajdu, nie po ich liczbie. Zmierzone
        # 2026-08-19: python-pptx liczy slajdy z `presentation.xml` i nigdy nie
        # zaglada do samej czesci slajdu, wiec slajd podmieniony na `<nonsense/>`
        # przechodzil. To ta sama klasa co „otwarcie to nie odczyt" przy
        # obrazach, gdzie dopiero porownanie pikseli nazwalo tolerancyjnych.
        slides, shapes = 0, 0
        for slide in deck.slides:
            slides += 1
            for shape in slide.shapes:
                shapes += 1
        if slides < 1:
            fail("python-pptx otworzyl prezentacje i nie znalazl w niej ani jednego slajdu")
        if shapes < 1:
            fail("python-pptx przeszedl po %d slajdzie(ach) i nie znalazl na nich ani jednego ksztaltu" % slides)
        return "python-pptx odczytal %d slajd(ow) i %d ksztalt(ow)" % (slides, shapes)

    return "brak niezaleznego czytnika"

def check_jpg(data):
    """JPEG walked segment by segment, against the specification.

    The tolerant readers are worth very little here and that is measured, not
    assumed (MVP-FORMATS.md 2.10). Of six deliberately broken files, GDI+ and
    exiftool called four of them fine - a flipped byte in the scan, bytes
    injected into the scan, a missing end marker and a file cut in half - and
    ffprobe accepted one whose signature was not even a JPEG signature. So
    this layer carries most of the weight for this format.

    What it asks that no reader here does: every segment length has to be
    consistent with where the next marker actually is, and nothing may follow
    the end of image marker.
    """
    if data[:2] != b"\xFF\xD8":
        fail("the file does not start with the start of image marker")
    if data[-2:] != b"\xFF\xD9":
        fail("the file does not end with the end of image marker")

    pos, segments, comments, comment_bytes = 2, 0, 0, 0
    scans, width, height = 0, None, None

    while pos + 1 < len(data):
        if data[pos] != 0xFF:
            fail(f"a marker was expected at offset {pos} and the file holds {data[pos]:#04x}")
        marker = data[pos + 1]

        if marker == 0xD8:
            fail(f"a second start of image marker at offset {pos}")
        if marker == 0xD9:
            pos += 2
            break
        if marker == 0x01 or 0xD0 <= marker <= 0xD7:
            # Markers that carry no payload at all.
            pos += 2
            continue

        if pos + 4 > len(data):
            fail(f"the segment at offset {pos} has no room for its length")
        length = struct.unpack(">H", data[pos + 2:pos + 4])[0]
        if length < 2:
            fail(f"the segment at offset {pos} declares a length of {length}, and the field counts itself")
        end = pos + 2 + length
        if end > len(data):
            fail(f"the segment at offset {pos} declares {length} B and the file ends after {len(data) - pos - 2}")

        segments += 1
        if marker == 0xFE:
            comments += 1
            comment_bytes += length - 2
        if marker in (0xC0, 0xC1, 0xC2, 0xC3) and length >= 7:
            height, width = struct.unpack(">HH", data[pos + 5:pos + 9])

        pos = end

        if marker == 0xDA:
            scans += 1
            # The entropy coded data has no length of its own. It ends at the
            # first marker that is neither a stuffed FF00 nor a restart.
            scan = pos
            while scan + 1 < len(data):
                if data[scan] != 0xFF:
                    scan += 1
                    continue
                nxt = data[scan + 1]
                if nxt == 0x00 or 0xD0 <= nxt <= 0xD7 or nxt == 0xFF:
                    scan += 2 if nxt != 0xFF else 1
                    continue
                break
            if scan + 1 >= len(data):
                fail("the scan runs to the end of the file without a closing marker")
            pos = scan

    if pos != len(data):
        fail(f"the end of image marker is at {pos} and {len(data) - pos} B follow it")
    if scans < 1:
        fail("there is no scan, so the file carries no image data")
    if width is None:
        fail("there is no frame header, so nothing says how large the picture is")
    ok(f"{width}x{height}, {segments} segments, {scans} scan(s), "
       f"{comments} comment(s) carrying {comment_bytes} B")


def check_tiff(data):
    # This layer earns more at TIFF than anywhere else, and that is measured
    # rather than assumed. On 2026-08-29 six deliberately broken files went
    # through five readers: a StripByteCounts announcing half the pixel data
    # was accepted by FOUR of them - Pillow, exiftool, WIC and GDI+ - and
    # caught only by x/image. A reader that decodes the image it can find does
    # not have to care that the directory lied about how much there was.
    if data[:2] != b"II":
        fail("the byte order mark is not II, and this generator writes little-endian only")
    if len(data) < 8:
        fail(f"the file is {len(data)} B and the header alone is 8 B")

    magic = struct.unpack("<H", data[2:4])[0]
    if magic != 42:
        fail(f"the magic number is {magic} and a TIFF says 42")

    ifd_at = struct.unpack("<I", data[4:8])[0]
    if ifd_at + 2 > len(data):
        fail(f"the directory is said to start at {ifd_at} and the file ends at {len(data)}")

    count = struct.unpack("<H", data[ifd_at:ifd_at + 2])[0]
    if count < 1:
        fail("the directory holds no entries")
    end = ifd_at + 2 + 12 * count + 4
    if end > len(data):
        fail(f"a directory of {count} entries runs to {end} and the file ends at {len(data)}")

    sizes = {1: 1, 2: 1, 3: 2, 4: 4, 5: 8, 7: 1}
    tags, previous = {}, -1
    for i in range(count):
        at = ifd_at + 2 + 12 * i
        tag, kind, n = struct.unpack("<HHI", data[at:at + 8])
        raw = data[at + 8:at + 12]
        # The specification requires ascending tag order and readers rely on
        # it. Nothing that merely decodes the picture would notice.
        if tag <= previous:
            fail(f"entry {i} has tag {tag} after tag {previous}, and a directory is in ascending order")
        previous = tag
        width = sizes.get(kind)
        if width is None:
            fail(f"entry {i} declares data type {kind}, which is not one this generator writes")
        total = width * n
        if total > 4:
            off = struct.unpack("<I", raw)[0]
            if off + total > len(data):
                fail(f"the value of tag {tag} runs to {off + total} and the file ends at {len(data)}")
            value = data[off:off + total]
        else:
            value = raw[:total]
        tags[tag] = (kind, n, value)

    nxt = struct.unpack("<I", data[ifd_at + 2 + 12 * count:end])[0]
    if nxt != 0:
        fail(f"the directory points at another one at {nxt}, and this generator writes a single page")

    def one(tag, name):
        if tag not in tags:
            fail(f"the directory has no {name}")
        kind, n, value = tags[tag]
        if n != 1:
            fail(f"{name} has {n} values and it should have one")
        return struct.unpack("<H", value)[0] if kind == 3 else struct.unpack("<I", value)[0]

    width = one(256, "ImageWidth")
    height = one(257, "ImageLength")
    compression = one(259, "Compression")
    photometric = one(262, "PhotometricInterpretation")
    offset = one(273, "StripOffsets")
    samples = one(277, "SamplesPerPixel")
    counts = one(279, "StripByteCounts")

    if width < 1 or height < 1:
        fail(f"the picture is {width}x{height}")
    if compression != 1:
        fail(f"the directory declares compression {compression} and this generator writes none")
    if photometric != 2:
        fail(f"the directory declares photometric {photometric} and this generator writes RGB")
    if samples != 3:
        fail(f"the directory declares {samples} samples per pixel and this generator writes three")

    if 258 not in tags:
        fail("the directory has no BitsPerSample")
    bits = tags[258][2]
    if len(bits) != samples * 2:
        fail(f"BitsPerSample carries {len(bits)} B for {samples} samples")
    for i in range(samples):
        depth = struct.unpack("<H", bits[i * 2:i * 2 + 2])[0]
        if depth != 8:
            fail(f"sample {i} is {depth} bits and this generator writes eight")

    # The defect four readers of five let through: the directory saying the
    # pixels are one length while the geometry says another.
    expected = width * height * samples
    if counts != expected:
        fail(f"the directory says the pixels are {counts} B and {width}x{height} at {samples} bytes is {expected} B")
    if offset < 8:
        fail(f"the pixels are said to start at {offset}, inside the header")
    if offset + counts > len(data):
        fail(f"the pixels run to {offset + counts} and the file ends at {len(data)}")
    # The pixels have to end where the directory begins, because this
    # generator puts its padding in front of them and nothing between.
    if offset + counts != ifd_at:
        fail(f"the pixels end at {offset + counts} and the directory begins at {ifd_at}")

    gap = offset - 8
    ok(f"{width}x{height}, {samples * 8} bit, uncompressed, gap {gap} B, "
       f"{count} entries in one directory, byte counts agree with the geometry")


def check_docx(data):
    opc_check(data, ["word/document.xml"], "Word document", "docx")


def check_xlsx(data):
    opc_check(data, ["xl/workbook.xml"], "Excel workbook", "xlsx")


def check_pptx(data):
    opc_check(data, ["ppt/presentation.xml",
                     "ppt/slideMasters/slideMaster1.xml",
                     "ppt/slideLayouts/slideLayout1.xml",
                     "ppt/theme/theme1.xml"], "PowerPoint presentation", "pptx")


def check_webp(data):
    if data[:4] != b"RIFF" or data[8:12] != b"WEBP":
        fail("the signature is not a WebP signature")
    if len(data) < 20:
        fail(f"the file is {len(data)} B and the container alone needs 20 B")

    declared = struct.unpack("<I", data[4:8])[0]
    # Everything the RIFF size covers, plus whatever sits after it. This
    # generator puts one to seven bytes there on purpose, because a chunk
    # block always costs an even number and half of every size would
    # otherwise be unreachable. More than seven means the padding went to the
    # wrong place, which no viewer would ever complain about.
    payload_end = 8 + declared
    if payload_end > len(data):
        fail(f"the header says the payload runs to {payload_end} and the file is {len(data)} B")
    trailing = len(data) - payload_end
    if trailing > 7:
        fail(f"{trailing} B sit after the RIFF payload, and only the odd byte belongs there - "
             "the bulk of the padding belongs in a chunk")

    pos, image, private, seen = 12, 0, 0, []
    while pos + 8 <= payload_end:
        tag = data[pos:pos + 4]
        size = struct.unpack("<I", data[pos + 4:pos + 8])[0]
        end = pos + 8 + size
        if end > payload_end:
            fail(f"the chunk {tag.decode('latin1')!r} says {size} B and only "
                 f"{payload_end - pos - 8} B are left")
        seen.append(tag.decode("latin1"))
        if tag in (b"VP8L", b"VP8 ", b"VP8X"):
            image += 1
            if tag == b"VP8L" and data[pos + 8] != 0x2F:
                fail("the lossless stream does not start with its signature byte")
        else:
            private += 1
        pos = end + (size & 1)

    if pos != payload_end:
        fail(f"the chunks end at {pos} and the payload ends at {payload_end}")
    if image != 1:
        fail(f"the file carries {image} image chunks and a still WebP has one")

    ok(f"riff {declared} B, chunks {seen}, {private} private, {trailing} B after the payload")


def iso_boxes(data, start, end, depth=0):
    """Walks one level of the box tree and yields (name, payload start, payload end).

    Shared by AVIF and JPEG XL, which both sit in ISO base media boxes.
    """
    pos = start
    while pos + 8 <= end:
        size = struct.unpack(">I", data[pos:pos + 4])[0]
        name = data[pos + 4:pos + 8]
        body = pos + 8
        if size == 1:
            if pos + 16 > end:
                fail("a box says its length is written in 64 bits and the file ends first")
            size = struct.unpack(">Q", data[pos + 8:pos + 16])[0]
            body = pos + 16
        elif size == 0:
            size = end - pos
        if size < body - pos:
            fail(f"the box {name.decode('latin1')!r} declares {size} B, which is shorter than its own header")
        stop = pos + size
        if stop > end:
            fail(f"the box {name.decode('latin1')!r} runs to {stop} and its parent ends at {end}")
        yield name, body, stop
        pos = stop
    if pos != end:
        fail(f"the boxes end at {pos} and the area they sit in ends at {end}")


def check_avif(data):
    if len(data) < 16:
        fail(f"the file is {len(data)} B and the smallest container needs more than that")

    top = list(iso_boxes(data, 0, len(data)))
    if not top:
        fail("the file carries no boxes at all")

    names = [n.decode("latin1") for n, _, _ in top]
    if names[0] != "ftyp":
        fail(f"the first box is {names[0]!r} and every ISO base media file starts with ftyp")

    _, ftyp_start, ftyp_end = top[0]
    if ftyp_end - ftyp_start < 8:
        fail("the ftyp box is too short to carry a brand")
    brands = [data[ftyp_start:ftyp_start + 4]]
    for at in range(ftyp_start + 8, ftyp_end, 4):
        brands.append(data[at:at + 4])
    if b"avif" not in brands:
        fail("no box says this file is an AVIF - the brand avif is in neither the major brand nor the compatible ones")

    if names.count("meta") != 1:
        fail(f"the file carries {names.count('meta')} meta boxes and an AVIF has one")
    if names.count("mdat") < 1:
        fail("the file carries no mdat box, so there is nothing for the picture to be stored in")

    # The metadata has to say what the item is and where it lives. Without
    # these a reader has a container and no picture, which is exactly the
    # shape a padding channel written to the wrong place would produce.
    meta_start, meta_end = next((s, e) for n, s, e in top if n == b"meta")
    inner = [n.decode("latin1") for n, _, _ in iso_boxes(data, meta_start + 4, meta_end)]
    for needed in ("hdlr", "pitm", "iloc", "iinf", "iprp"):
        if needed not in inner:
            fail(f"the meta box has no {needed} box, so it does not say {AVIF_MEANING[needed]}")

    # Padding belongs after the picture and nowhere else.
    free_bytes = 0
    seen_mdat = False
    for name, start, stop in top:
        if name == b"mdat":
            seen_mdat = True
        elif name == b"free":
            free_bytes += stop - start + 8
            if not seen_mdat:
                fail("a free box sits before the picture data, and this generator puts padding after it")

    ok(f"boxes {names}, brands {[b.decode('latin1') for b in brands]}, {free_bytes} B of free")


JXL_SIGNATURE = b"\x00\x00\x00\x0cJXL \x0d\x0a\x87\x0a"


def check_jxl(data):
    """A JPEG XL in the container form this generator writes.

    Two shapes are legal for the format - a bare codestream starting FF 0A, and
    the container. This tool always writes the container, because that is where
    a free box can live, so a bare codestream here means the writer changed and
    the padding went somewhere undefined.
    """
    if len(data) < len(JXL_SIGNATURE) + 8:
        fail(f"the file is {len(data)} B and the smallest container needs more than that")

    if data[:2] == b"\xff\x0a":
        fail("the file is a bare codestream, and this generator writes the container form")
    if data[:len(JXL_SIGNATURE)] != JXL_SIGNATURE:
        fail("the file does not open with the JPEG XL signature box")

    top = list(iso_boxes(data, 0, len(data)))
    names = [n.decode("latin1") for n, _, _ in top]
    if names[0] != "JXL ":
        fail(f"the first box is {names[0]!r} and a JPEG XL container starts with the signature box")
    if len(names) < 2 or names[1] != "ftyp":
        fail("the signature box is not followed by ftyp, so nothing says what this file is")

    _, ftyp_start, ftyp_end = top[1]
    if ftyp_end - ftyp_start < 8:
        fail("the ftyp box is too short to carry a brand")
    brands = [data[ftyp_start:ftyp_start + 4]]
    for at in range(ftyp_start + 8, ftyp_end, 4):
        brands.append(data[at:at + 4])
    if brands[0] != b"jxl ":
        fail(f"the major brand is {brands[0]!r} and a JPEG XL says jxl ")
    # Measured 2026-08-31: libjxl refuses the file when the compatible brand is
    # missing, while the pure Go decoder accepts it. The stricter reader is the
    # one worth writing for, so this checker asks for it too.
    if b"jxl " not in brands[1:]:
        fail("ftyp names no compatible brand jxl , which libjxl refuses")

    if names.count("jxlc") != 1:
        fail(f"the file carries {names.count('jxlc')} jxlc boxes and this generator writes one")

    _, code_start, code_end = next((n, s, e) for n, s, e in top if n == b"jxlc")
    if data[code_start:code_start + 2] != b"\xff\x0a":
        fail("the jxlc box does not start with the codestream signature FF 0A")

    # Padding belongs after the picture and nowhere else.
    free_bytes = 0
    seen_code = False
    for name, start, stop in top:
        if name == b"jxlc":
            seen_code = True
        elif name == b"free":
            free_bytes += stop - start + 8
            if not seen_code:
                fail("a free box sits before the picture data, and this generator puts padding after it")

    tail = len(data) - top[-1][2]
    if tail:
        fail(f"{tail} B sit after the last box, outside any box at all")

    ok(f"boxes {names}, brands {[b.decode('latin1') for b in brands]}, "
       f"codestream {code_end - code_start} B, {free_bytes} B of free")


AVIF_MEANING = {
    "hdlr": "what kind of thing it holds",
    "pitm": "which item is the picture",
    "iloc": "where the picture data is",
    "iinf": "what the item is",
    "iprp": "how wide and tall the picture is",
}


CHECKS = {"png": check_png, "wav": check_wav, "pdf": check_pdf, "zip": check_zip,
          "log": check_log, "csv": check_csv, "json": check_json, "xml": check_xml,
          "svg": check_svg, "html": check_html, "targz": check_targz,
          "bmp": check_bmp, "gif": check_gif, "ico": check_ico, "jpg": check_jpg,
          "tiff": check_tiff, "webp": check_webp, "avif": check_avif, "jxl": check_jxl,
          "docx": check_docx, "xlsx": check_xlsx, "pptx": check_pptx}

# Checks that take the shape of the file as well as its bytes. Everything else
# is handed the bytes alone, so adding a setting to one check cannot change how
# any other one is called.
TAKES_SETTINGS = {"csv"}

if __name__ == "__main__":
    if len(sys.argv) < 3 or sys.argv[1] not in CHECKS:
        print("FAIL usage: strict.py <" + "|".join(CHECKS) + "> <path> [key=value ...]")
        sys.exit(1)
    kind = sys.argv[1]
    extra = {}
    for word in sys.argv[3:]:
        if "=" not in word:
            print(f"FAIL setting {word!r} is not key=value")
            sys.exit(1)
        key, _, value = word.partition("=")
        extra[key] = value
    if extra and kind not in TAKES_SETTINGS:
        print(f"FAIL the {kind} check takes no settings, and was given {sorted(extra)}")
        sys.exit(1)
    with open(sys.argv[2], "rb") as handle:
        body = handle.read()
    if kind in TAKES_SETTINGS:
        CHECKS[kind](body, extra)
    else:
        CHECKS[kind](body)
